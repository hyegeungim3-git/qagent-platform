# -*- coding: utf-8 -*-
"""
gen_spec_pptx.py — 화면설계서 PPTX 생성기 (범용)
captures/manifest.json + captures/<viewport>/*.png → 화면설계서 PPTX.
화면별 슬라이드: 스크린샷(번호 원형 마커) + Description 표(개요·경로·기능 명세·CRUD).
사용: python gen_spec_pptx.py [캡처폴더=captures] [출력.pptx]
"""
import json, os, sys, datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
CAP = os.path.join(HERE, sys.argv[1] if len(sys.argv) > 1 else 'captures')
CROP = os.path.join(CAP, '_pptx_crops'); os.makedirs(CROP, exist_ok=True)

# 한글 폰트 — 한국어 패밀리명으로 지정 (weight별 분리 패밀리)
FONT_R = '나눔스퀘어 네오 Regular'
FONT_B = '나눔스퀘어 네오 Bold'
FONT_EB = '나눔스퀘어 네오 ExtraBold'

NAVY = RGBColor(0x16, 0x21, 0x3F)
PRIMARY = RGBColor(0x35, 0x66, 0xCD)
MUTED = RGBColor(0x58, 0x61, 0x83)
FAINT = RGBColor(0x8A, 0x92, 0xB0)
LINE = RGBColor(0xD9, 0xE0, 0xF0)
BG_SOFT = RGBColor(0xF3, 0xF6, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

with open(os.path.join(CAP, 'manifest.json'), encoding='utf-8') as f:
    M = json.load(f)
TODAY = datetime.date.today().strftime('%Y.%m.%d')
OUT = sys.argv[2] if len(sys.argv) > 2 else os.path.join(
    HERE, '..', '화면설계서_' + datetime.date.today().strftime('%Y%m%d') + '.pptx')
DOC_TITLE = M['target'].get('docTitle', '화면설계서')
DOC_SUB = M['target'].get('docSub', M['target']['name'])


def set_font(run, size, bold=False, color=NAVY, mono=False, xb=False):
    """weight별 한국어 패밀리(latin+eastAsia 동시 지정)."""
    name = 'Consolas' if mono else (FONT_EB if xb else (FONT_B if bold else FONT_R))
    f = run.font
    f.name = name
    f.size = Pt(size); f.bold = bold; f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT_EB if xb else (FONT_B if bold else FONT_R))


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def para(tf, text, size, bold=False, color=NAVY, before=0, first=False, mono=False, xb=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    if before: p.space_before = Pt(before)
    r = p.add_run(); r.text = text; set_font(r, size, bold, color, mono, xb)
    return p


def cell_borders(cell, color='C9D3E8', w=9525):
    """표 셀 4변 테두리 — 정식 설계서 표 형태."""
    tcPr = cell._tc.get_or_add_tcPr()
    for i, tag in enumerate(('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB')):
        ln = tcPr.find(qn(tag))
        if ln is None:
            ln = tcPr.makeelement(qn(tag), {})
            tcPr.insert(i, ln)
        ln.set('w', str(w)); ln.set('cap', 'flat')
        for ch in list(ln):
            ln.remove(ch)
        fill = ln.makeelement(qn('a:solidFill'), {}); ln.append(fill)
        clr = fill.makeelement(qn('a:srgbClr'), {'val': color}); fill.append(clr)


def crop_top(png, max_px):
    im = Image.open(png)
    if im.height <= max_px:
        return png, im.width, im.height
    out = os.path.join(CROP, os.path.basename(os.path.dirname(png)) + '-' + str(max_px) +
                       '-' + os.path.basename(png))
    # 원본이 더 새로우면 캐시 무효 (재캡처 후 옛 크롭이 문서에 남는 결함 방지)
    if (not os.path.exists(out)) or os.path.getmtime(out) < os.path.getmtime(png):
        im.crop((0, 0, im.width, max_px)).save(out)
    return out, im.width, max_px


def fit(w_px, h_px, max_w_in, max_h_in):
    r = min(max_w_in / w_px, max_h_in / h_px)
    return w_px * r, h_px * r


def footer(slide, page_no):
    tf = box(slide, 0.5, 7.12, 9.0, 0.3).text_frame
    para(tf, DOC_TITLE + '  ·  ' + DOC_SUB, 7.5, False, FAINT, first=True)
    tf2 = box(slide, 12.35, 7.12, 0.5, 0.3).text_frame
    p = para(tf2, str(page_no), 8, False, FAINT, first=True)
    p.alignment = PP_ALIGN.RIGHT


prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
page = 0

# ---------- 표지 ----------
page += 1
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.18))
band.fill.solid(); band.fill.fore_color.rgb = PRIMARY; band.line.fill.background()
tf = box(s, 1.0, 2.3, 11.3, 2.6).text_frame
para(tf, DOC_TITLE, 40, True, NAVY, first=True, xb=True)
para(tf, DOC_SUB, 20, False, PRIMARY, before=8)
para(tf, 'Interface Design Specification', 13, False, FAINT, before=4)
tf2 = box(s, 1.0, 5.6, 11.3, 1.2).text_frame
para(tf2, '기획 아스트로젠 · 개발 오큐브    |    기준 ' + M['target']['name'].split('—')[0].strip()
     + ' ' + M['target'].get('version', ''), 12, False, MUTED, first=True)
para(tf2, '작성일 ' + TODAY + '  ·  본 문서의 전 화면 스크린샷과 기능 마커 좌표는 실행 중인 프로토타입에서 자동 생성된 실측값입니다.',
     12, False, MUTED, before=3)

# ---------- 문서 개요 ----------
page += 1
s = prs.slides.add_slide(BLANK)
tf = box(s, 0.6, 0.4, 12.1, 0.7).text_frame
para(tf, '문서 개요', 24, True, first=True, xb=True)
tf = box(s, 0.6, 1.3, 5.9, 5.4).text_frame
para(tf, '화면 ID 체계', 15, True, PRIMARY, first=True)
sc = M['viewports'].get('desktop', [])
seen, id_sys = set(), []
for sc_i in sc:
    pfx = sc_i['id'].split('-')[0]
    if pfx not in seen:
        seen.add(pfx); id_sys.append((pfx, sc_i['cat']))
for code, mean in id_sys:
    para(tf, '·  ' + code + ' — ' + mean, 12, False, MUTED, before=4)
para(tf, '표기 규칙', 15, True, PRIMARY, before=14)
for t in ['화면 스크린샷 위 번호 ●는 우측 Description 표의 기능 번호와 1:1 대응',
          '번호 좌표는 캡처 시 실제 UI 요소 위치를 자동 측정한 값',
          'CRUD = 해당 화면에서 일어나는 데이터 생성(C)·조회(R)·수정(U)·삭제(D)']:
    para(tf, '·  ' + t, 11, False, MUTED, before=4)
tf = box(s, 6.9, 1.3, 5.8, 5.4).text_frame
para(tf, '작성 체계 (재현 가능)', 15, True, PRIMARY, first=True)
for t in ['화면 정의(경로·설명·기능 명세·CRUD·마커 셀렉터)는 매니페스트 단일 정본으로 관리',
          '캡처는 헤드리스 크롬이 매니페스트를 순회하며 자동 수행 (데스크톱 1280 · 모바일 375)',
          '본 문서는 캡처 산출물에서 자동 생성 — 프로토타입 변경 시 재실행으로 최신화',
          '재생성: screen-spec-tool에서  node capture.mjs all  →  python gen_spec_pptx.py']:
    para(tf, '·  ' + t, 12, False, MUTED, before=6)
para(tf, ' ', 8, before=6)
para(tf, '수록: 화면 ' + str(len(sc)) + '종 × 데스크톱 풀페이지 + 모바일 부록 (총 '
     + str(sum(len(v) for v in M['viewports'].values())) + '컷)', 12, True, NAVY, before=8)
footer(s, page)

# ---------- 화면 목록 ----------
PER = 28
for pg in range(0, len(sc), PER):
    page += 1
    chunk = sc[pg:pg + PER]
    s = prs.slides.add_slide(BLANK)
    tf = box(s, 0.6, 0.35, 12.1, 0.6).text_frame
    ttl = '화면 목록' + ('' if len(sc) <= PER else f' ({pg // PER + 1}/{-(-len(sc) // PER)})')
    para(tf, ttl, 24, True, first=True, xb=True)
    tbl = s.shapes.add_table(len(chunk) + 1, 4, Inches(0.6), Inches(1.05), Inches(12.1), Inches(5.9)).table
    tbl.columns[0].width = Inches(1.1); tbl.columns[1].width = Inches(1.4)
    tbl.columns[2].width = Inches(3.6); tbl.columns[3].width = Inches(6.0)
    for c, h in enumerate(['ID', '분류', '화면명', '경로 · 접근']):
        cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        r = cell.text_frame.paragraphs[0].add_run(); r.text = h
        set_font(r, 10.5, True, WHITE)
        cell_borders(cell)
    for i, sc_i in enumerate(chunk):
        vals = [sc_i['id'], sc_i['cat'], sc_i['title'], sc_i['route'] + '  ·  ' + sc_i['nav']]
        for c, v in enumerate(vals):
            cell = tbl.cell(i + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SOFT if i % 2 else WHITE
            cell.margin_top = Emu(9525); cell.margin_bottom = Emu(9525)
            r = cell.text_frame.paragraphs[0].add_run(); r.text = v
            set_font(r, 9.5, c == 0, NAVY if c < 3 else MUTED, mono=(c == 0))
            cell_borders(cell)
    footer(s, page)

# ---------- 화면별 슬라이드 ----------
AUTH_KO = { 'public': '공개(비로그인)', 'parent': '보호자(로그인)', 'admin': '관리자' }


def fill_cell(cell, bg=WHITE, anchor=MSO_ANCHOR.TOP):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Emu(54864); cell.margin_right = Emu(54864)
    cell.margin_top = Emu(27432); cell.margin_bottom = Emu(27432)
    cell.vertical_anchor = anchor
    cell_borders(cell)


def label_cell(cell, text):
    fill_cell(cell, BG_SOFT, MSO_ANCHOR.MIDDLE)
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; set_font(r, 9.5, True, NAVY)


def value_cell(cell, text, mono=False, color=NAVY, size=9.5, bold=False):
    fill_cell(cell, WHITE, MSO_ANCHOR.MIDDLE)
    p = cell.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text; set_font(r, size, bold, color, mono=mono)


def est_lines(text, cpl):
    return max(1, -(-len(text or '') // cpl))


for i, sc_i in enumerate(sc):
    page += 1
    s = prs.slides.add_slide(BLANK)
    feats = sc_i.get('features') or []
    crud = sc_i.get('crud')

    # ---------- ① 상단 메타데이터 표 (화면 ID·화면명·분류 / 경로·접근·권한 / 화면 개요) ----------
    MT_X, MT_Y, MT_W = 0.5, 0.32, 12.35
    mt = s.shapes.add_table(3, 6, Inches(MT_X), Inches(MT_Y), Inches(MT_W), Inches(1.3)).table
    mt.first_row = False; mt.horz_banding = False
    widths = [1.0, 2.2, 1.0, 4.55, 1.0, 2.6]
    for c, w in enumerate(widths):
        mt.columns[c].width = Inches(w)
    label_cell(mt.cell(0, 0), '화면 ID')
    value_cell(mt.cell(0, 1), sc_i['id'], mono=True, color=PRIMARY, size=10.5, bold=True)
    label_cell(mt.cell(0, 2), '화면명')
    value_cell(mt.cell(0, 3), sc_i['title'], size=10.5, bold=True)
    label_cell(mt.cell(0, 4), '분류')
    value_cell(mt.cell(0, 5), sc_i['cat'])
    label_cell(mt.cell(1, 0), '경로')
    value_cell(mt.cell(1, 1), sc_i['route'], mono=True, color=PRIMARY)
    label_cell(mt.cell(1, 2), '접근 경로')
    value_cell(mt.cell(1, 3), sc_i['nav'])
    label_cell(mt.cell(1, 4), '접근 권한')
    value_cell(mt.cell(1, 5), AUTH_KO.get(sc_i.get('auth'), '-'))
    label_cell(mt.cell(2, 0), '화면 개요')
    dc = mt.cell(2, 1); dc.merge(mt.cell(2, 5))
    fill_cell(dc, WHITE, MSO_ANCHOR.MIDDLE)
    p = dc.text_frame.paragraphs[0]
    r = p.add_run(); r.text = sc_i['desc']; set_font(r, 9.5, False, MUTED)
    for ri, h in [(0, 0.3), (1, 0.3), (2, 0.42)]:
        mt.rows[ri].height = Inches(h)

    desc_lines = est_lines(sc_i['desc'], 78)
    CONTENT_Y = MT_Y + 0.3 + 0.3 + max(0.42, 0.18 + desc_lines * 0.16) + 0.14

    # ---------- ② 스크린샷 (좌) + 번호 마커 ----------
    # 크롭 높이는 '상한(3400px) 안에 실제로 들어오는 가장 깊은 마커 + 여백'까지만 확장
    # (상한 밖 초심층 마커 때문에 이미지만 얇아지는 것을 방지)
    IMG_X, IMG_Y = 0.5, CONTENT_Y
    png = os.path.join(CAP, 'desktop', sc_i['file'])
    _dsf = sc_i.get('dsf', 1)
    _ys = [m['y'] * _dsf for m in (sc_i.get('markers') or [])]
    _inside = [y for y in _ys if y <= 3180]
    crop_px = int(min(3400, max(2000, (max(_inside) + 220) if _inside else 2000)))
    use, w_px, h_px = crop_top(png, crop_px)
    w_in, h_in = fit(w_px, h_px, 6.55, 7.02 - IMG_Y - 0.22)
    pic = s.shapes.add_picture(use, Inches(IMG_X), Inches(IMG_Y), Inches(w_in), Inches(h_in))
    pic.line.color.rgb = LINE; pic.line.width = Pt(0.75)
    if use != png:
        skipped = sorted(m['n'] for m in (sc_i.get('markers') or []) if m['y'] * _dsf > h_px - 8)
        note = ('※ 상단 일부 — 전체 화면은 ' + os.path.basename(CAP) + '/desktop/' + sc_i['file'] + ' 참조'
                + ((' · 번호 ' + '·'.join(map(str, skipped)) + '는 화면 하단 영역(원본 참조)') if skipped else ''))
        cap_tf = box(s, IMG_X, IMG_Y + h_in + 0.03, 6.75, 0.28).text_frame
        para(cap_tf, note, 8, False, FAINT, first=True)
    scale = w_in / w_px
    dsf = sc_i.get('dsf', 1)
    for m in (sc_i.get('markers') or []):
        mx, my = m['x'] * dsf, m['y'] * dsf
        if my > h_px - 8:
            continue
        d = 0.26
        cx = max(IMG_X, min(IMG_X + w_in - d, IMG_X + mx * scale - d / 2))
        cy = max(IMG_Y, min(IMG_Y + h_in - d, IMG_Y + my * scale - d / 2))
        ov = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
        ov.fill.solid(); ov.fill.fore_color.rgb = PRIMARY
        ov.line.color.rgb = WHITE; ov.line.width = Pt(1.25)
        ov.shadow.inherit = False
        otf = ov.text_frame
        otf.margin_left = otf.margin_right = otf.margin_top = otf.margin_bottom = 0
        otf.vertical_anchor = MSO_ANCHOR.MIDDLE
        op = otf.paragraphs[0]; op.alignment = PP_ALIGN.CENTER
        orun = op.add_run(); orun.text = str(m['n'])
        set_font(orun, 11, True, WHITE)

    # ---------- ③ Description 표 (우) — 번호 기능 명세 ----------
    TX, TW = 7.3, 5.55
    n_rows = 1 + len(feats)
    tbl = s.shapes.add_table(n_rows, 2, Inches(TX), Inches(CONTENT_Y), Inches(TW), Inches(0.34 * n_rows)).table
    tbl.first_row = False; tbl.horz_banding = False
    tbl.columns[0].width = Inches(0.5)
    tbl.columns[1].width = Inches(TW - 0.5)
    hc = tbl.cell(0, 0); hc.merge(tbl.cell(0, 1))
    hc.fill.solid(); hc.fill.fore_color.rgb = NAVY
    hc.vertical_anchor = MSO_ANCHOR.MIDDLE
    hp = hc.text_frame.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run(); hr.text = 'Description'; set_font(hr, 11.5, True, WHITE)
    cell_borders(hc)
    tbl.rows[0].height = Inches(0.32)
    desc_h = 0.32
    for fi, ft_ in enumerate(feats):
        ri = 1 + fi
        label = ft_ if isinstance(ft_, str) else ft_['label']
        body = None if isinstance(ft_, str) else ft_.get('desc')
        c0 = tbl.cell(ri, 0); fill_cell(c0)
        p0 = c0.text_frame.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = str(fi + 1); set_font(r0, 10.5, True, PRIMARY)
        c1 = tbl.cell(ri, 1); fill_cell(c1)
        p1 = c1.text_frame.paragraphs[0]
        r1 = p1.add_run(); r1.text = label; set_font(r1, 9.5, True, NAVY)
        if body:
            pb = c1.text_frame.add_paragraph(); pb.space_before = Pt(1)
            rb = pb.add_run(); rb.text = body; set_font(rb, 8.5, False, MUTED)
        desc_h += 0.20 + est_lines(body, 46) * 0.145 if body else 0.30

    # ---------- ④ CRUD 표 (우측 하단, 별도) ----------
    if crud:
        crud_h = 0.30 + 4 * 0.27
        cy0 = max(CONTENT_Y + desc_h + 0.18, 7.06 - crud_h)
        cy0 = min(cy0, 7.45 - crud_h)
        ct = s.shapes.add_table(5, 2, Inches(TX), Inches(cy0), Inches(TW), Inches(crud_h)).table
        ct.first_row = False; ct.horz_banding = False
        ct.columns[0].width = Inches(0.9)
        ct.columns[1].width = Inches(TW - 0.9)
        hc = ct.cell(0, 0); hc.merge(ct.cell(0, 1))
        hc.fill.solid(); hc.fill.fore_color.rgb = NAVY
        hc.vertical_anchor = MSO_ANCHOR.MIDDLE
        hp = hc.text_frame.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run(); hr.text = 'CRUD'; set_font(hr, 11.5, True, WHITE)
        cell_borders(hc)
        ct.rows[0].height = Inches(0.28)
        for ki, (k, kn) in enumerate([('C', '생성'), ('R', '조회'), ('U', '수정'), ('D', '삭제')]):
            ri = ki + 1
            c0 = ct.cell(ri, 0); fill_cell(c0, BG_SOFT, MSO_ANCHOR.MIDDLE)
            p0 = c0.text_frame.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
            r0 = p0.add_run(); r0.text = k; set_font(r0, 10, True, PRIMARY, mono=True)
            r0b = p0.add_run(); r0b.text = ' ' + kn; set_font(r0b, 8.5, True, NAVY)
            c1 = ct.cell(ri, 1); fill_cell(c1, WHITE, MSO_ANCHOR.MIDDLE)
            p1 = c1.text_frame.paragraphs[0]
            r1 = p1.add_run(); r1.text = crud.get(k) or '—'
            set_font(r1, 9, False, MUTED if (crud.get(k) or '—') == '—' else NAVY)
            ct.rows[ri].height = Inches(0.24)
    footer(s, page)

# ---------- 모바일 부록 ----------
mob = M['viewports'].get('mobile', [])
if mob:
    per = 4
    for gi in range(0, len(mob), per):
        page += 1
        group = mob[gi:gi + per]
        s = prs.slides.add_slide(BLANK)
        tf = box(s, 0.6, 0.3, 12.1, 0.6).text_frame
        para(tf, '모바일(375px) 부록 — ' + group[0]['id'] + ' ~ ' + group[-1]['id'], 18, True, first=True, xb=True)
        for k, sc_i in enumerate(group):
            png = os.path.join(CAP, 'mobile', sc_i['file'])
            use, w_px, h_px = crop_top(png, 2800)
            w_in, h_in = fit(w_px, h_px, 2.85, 5.4)
            x = 0.6 + k * 3.15
            pic = s.shapes.add_picture(use, Inches(x), Inches(1.12), Inches(w_in), Inches(h_in))
            pic.line.color.rgb = LINE; pic.line.width = Pt(0.75)
            cap_tf = box(s, x, 1.12 + h_in + 0.04, 3.0, 0.4).text_frame
            para(cap_tf, sc_i['id'] + ' ' + sc_i['title'], 9.5, True, MUTED, first=True)
        footer(s, page)

prs.save(OUT)
print('저장:', os.path.abspath(OUT))
print('슬라이드:', len(prs.slides._sldIdLst))
